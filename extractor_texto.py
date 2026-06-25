"""
extractor_texto.py — Extrae texto plano de PDF, DOCX, XLSX, PPTX, TXT.
Sin dependencias de GUI. Devuelve el texto o lanza excepción si no puede.

Sprint E (issue #26): para PDFs escaneados (sin texto seleccionable),
hay una capa de OCR de respaldo vía tesseract + poppler. Requiere los
binarios de sistema 'tesseract-ocr' y 'poppler-utils' instalados (no
solo las librerías de Python pytesseract/pdf2image) — si no están
disponibles, el OCR simplemente no se intenta y el resto del pipeline
sigue funcionando igual que antes de este sprint (el campo queda en
'pendiente_revision', no se inventa nada).
"""
from __future__ import annotations
import os, re

MAX_CHARS = 8000   # límite de texto a extraer (suficiente para un resumen)

# Si pdfplumber extrae menos que esto, se asume que el PDF es (o puede
# ser) escaneado y se intenta OCR como respaldo.
UMBRAL_TEXTO_ESCANEADO = 20

# Páginas máximas a pasar por OCR — mucho más caro por página que
# pdfplumber, así que el límite es más bajo que el de extracción normal.
MAX_PAGINAS_OCR = 5

OCR_DISPONIBLE = None  # se calcula una sola vez, lazy


def _ocr_disponible() -> bool:
    """Detecta si los binarios de sistema (no solo las libs de Python)
    están instalados. Cacheado: solo se chequea una vez por proceso."""
    global OCR_DISPONIBLE
    if OCR_DISPONIBLE is not None:
        return OCR_DISPONIBLE
    try:
        import shutil
        import pytesseract  # noqa: F401
        import pdf2image  # noqa: F401
        OCR_DISPONIBLE = bool(shutil.which("tesseract") and shutil.which("pdftoppm"))
    except ImportError:
        OCR_DISPONIBLE = False
    return OCR_DISPONIBLE


def extraer(path: str) -> str:
    """
    Extrae texto del archivo indicado.
    Devuelve string (puede estar vacío si el archivo no tiene texto seleccionable).
    Lanza FileNotFoundError si el archivo no existe.
    Lanza ValueError si la extensión no está soportada.
    """
    return extraer_con_fuente(path)[0]


def extraer_con_fuente(path: str) -> tuple[str, str]:
    """
    Igual que extraer(), pero además devuelve la fuente del texto:
    'texto_directo' (capa nativa del archivo) u 'ocr' (PDF escaneado,
    resuelto vía tesseract). Para formatos que no son PDF, la fuente
    siempre es 'texto_directo' — no hay heurística de "escaneado" para
    DOCX/XLSX/PPTX/TXT, que siempre tienen texto seleccionable por diseño.
    """
    if not os.path.isfile(path):
        raise FileNotFoundError(f"Archivo no encontrado: {path}")

    ext = os.path.splitext(path)[1].lower()

    if ext == ".pdf":
        return _pdf_con_fuente(path)
    elif ext in (".docx", ".doc"):
        return _docx(path), "texto_directo"
    elif ext in (".xlsx", ".xls"):
        return _xlsx(path), "texto_directo"
    elif ext == ".pptx":
        return _pptx(path), "texto_directo"
    elif ext == ".txt":
        return _txt(path), "texto_directo"
    else:
        raise ValueError(f"Extensión no soportada: {ext}")


# ── Extractores por formato ───────────────────────────────────────────────────

def _pdf_con_fuente(path: str) -> tuple[str, str]:
    import pdfplumber
    texts = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages[:15]:   # primeras 15 páginas es más que suficiente
            t = page.extract_text()
            if t:
                texts.append(t)
    texto = _clean("\n".join(texts))

    if len(texto) >= UMBRAL_TEXTO_ESCANEADO or not _ocr_disponible():
        return texto, "texto_directo"

    # Texto insuficiente y OCR disponible: probablemente un PDF escaneado.
    texto_ocr = _ocr_pdf(path)
    if len(texto_ocr) > len(texto):
        return texto_ocr, "ocr"
    return texto, "texto_directo"  # OCR no mejoró nada, no inventar


def _ocr_pdf(path: str) -> str:
    """OCR vía pdf2image (rasteriza con poppler) + pytesseract.
    Tolerante a fallos: cualquier excepción de la cadena de binarios
    devuelve '' en vez de propagar — un PDF que no se puede OCRizar no
    debe romper el resto del pipeline, queda en pendiente_revision.
    Pero SÍ se loguea como warning (issue #20): un fallo silencioso acá
    puede esconder un problema de instalación real (ej. tesseract
    instalado pero sin el paquete de idioma 'spa'), que de otro modo el
    usuario nunca se entera que existe — solo vería "el sistema no
    extrae datos de PDFs escaneados" y asumiría que es un límite del
    producto, no un problema de infraestructura corregible."""
    try:
        import pytesseract
        from pdf2image import convert_from_path
        paginas = convert_from_path(path, dpi=200, first_page=1,
                                    last_page=MAX_PAGINAS_OCR)
        textos = []
        for img in paginas:
            t = pytesseract.image_to_string(img, lang="spa+eng")
            if t:
                textos.append(t)
        return _clean("\n".join(textos))
    except Exception as e:
        import logging
        logging.warning(f"OCR falló para {path!r}: {e}")
        return ""


def _pdf(path: str) -> str:
    """Mantenido por compatibilidad — usar extraer() o extraer_con_fuente()."""
    return _pdf_con_fuente(path)[0]


def _docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())
    # Incluir tablas
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
            if row_text:
                parts.append(row_text)
    return _clean("\n".join(parts))


def _xlsx(path: str) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets[:3]:   # primeras 3 hojas
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join(
                str(c).strip() for c in row if c is not None and str(c).strip()
            )
            if row_text:
                parts.append(row_text)
    return _clean("\n".join(parts))


def _pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
    return _clean("\n".join(parts))


def _txt(path: str) -> str:
    for enc in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
        try:
            with open(path, "r", encoding=enc) as f:
                return _clean(f.read())
        except (UnicodeDecodeError, LookupError):
            continue
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return _clean(f.read())


# ── Helpers ───────────────────────────────────────────────────────────────────

def _clean(text: str) -> str:
    """Normaliza whitespace y trunca al límite de caracteres."""
    # Colapsar líneas en blanco múltiples
    text = re.sub(r"\n{3,}", "\n\n", text)
    # Colapsar espacios múltiples en la misma línea
    text = re.sub(r"[ \t]{2,}", " ", text)
    return text.strip()[:MAX_CHARS]
